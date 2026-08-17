from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import os

def main():
    prs = Presentation()

    # Brand Colors
    NAVY = RGBColor(0x00, 0x41, 0x79)
    GOLD = RGBColor(0xFF, 0xCF, 0x00)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    STONE = RGBColor(0xE8, 0xEE, 0xF4)
    NAVY_LIGHT = RGBColor(0x1A, 0x5A, 0x94)

    # Specific Logo Path requested by User
    logo_path = os.path.join("public", "logo A png.png")

    def apply_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = NAVY

    def apply_brand_to_slide(slide):
        apply_background(slide)
        if os.path.exists(logo_path):
            left = Inches(8.5)
            top = Inches(0.2)
            height = Inches(0.9)
            slide.shapes.add_picture(logo_path, left, top, height=height)
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(1.6), Inches(8), Pt(3)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = GOLD
        line.line.color.rgb = GOLD

    # --- Title Slide ---
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    apply_background(slide)
    title = slide.shapes.title
    title.text = "የመዝሙር ክፍል አደረጃጀትና የአገልግሎት መዋቅር ማሻሻያ ጥናታዊ የውሳኔ ሐሳብ"
    title.top = Inches(2.3)
    for paragraph in title.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.color.rgb = GOLD
            run.font.bold = True
            run.font.name = "Segoe UI"
            run.font.size = Pt(40)
    subtitle = slide.shapes.placeholders[1]
    subtitle.text = "(ሙሉ ሰነድ)"
    subtitle.top = Inches(4.5)
    for paragraph in subtitle.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.color.rgb = STONE
            run.font.bold = False
            run.font.name = "Segoe UI"
            run.font.size = Pt(28)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(3.5), Inches(5.2), Inches(3), Pt(5)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.color.rgb = GOLD
    if os.path.exists(logo_path):
        left = Inches(4.25)
        top = Inches(0.5)
        height = Inches(1.5)
        slide.shapes.add_picture(logo_path, left, top, height=height)

    # Helper for adding regular slides
    def add_slide(title_text, content_list, font_size=20):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title_shape = slide.shapes.title
        title_shape.text = title_text
        title_shape.left = Inches(0.5)
        title_shape.top = Inches(0.3)
        title_shape.width = Inches(7.5)
        title_shape.height = Inches(1.2)
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.LEFT
            for run in paragraph.runs:
                run.font.color.rgb = GOLD
                run.font.bold = True
                run.font.name = "Segoe UI"
                run.font.size = Pt(30)
        tf = slide.shapes.placeholders[1].text_frame
        tf.margin_top = Inches(0.2)
        slide.shapes.placeholders[1].top = Inches(1.8)
        slide.shapes.placeholders[1].left = Inches(0.5)
        slide.shapes.placeholders[1].width = Inches(9.0)
        slide.shapes.placeholders[1].height = Inches(5.4)
        for i, content in enumerate(content_list):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = content
            p.level = 0
            p.space_after = Pt(14)
            for run in p.runs:
                if ":" in content and i == 0:
                    run.font.color.rgb = WHITE
                    run.font.bold = True
                else:
                    run.font.color.rgb = STONE
                run.font.name = "Segoe UI"
                run.font.size = Pt(font_size)
        apply_brand_to_slide(slide)

    # ── SLIDE 2: ምስጋና ──────────────────────────────────────────────
    add_slide("ምስጋና (Acknowledgment)", [
        "በቅድመ አያቶቻችን እምነትና ተጋድሎ የተመሰረተችው ቅድስት ቤተክርስቲያናችን፣ ከትውልድ ወደ ትውልድ ያሬዳዊ ዜማዋንና መንፈሳዊ ሀብቷን ጠብቃ እንድታቆይልን ላደረገ ለአምላካችን እግዚአብሔር ምስጋና ይሁን።",
        "ሰንበት ትምህርት ቤታችን ከተመሰረተበት 1965 ዓ.ም ጀምሮ እስካሁን ድረስ የመዝሙር አገልግሎቱ ሳይቋረጥና ያሬዳዊ ዜማውን ሳይለቅ አሁን ላለበት የዕድገት ደረጃ እንዲደርስ ሌሊትና ቀን በጸሎት፣ በጉልበት፣ በዕውቀትና በገንዘባቸው ላገለገሉ ቀደምት አባቶች፣ እናቶች፣ መምህራንና የመዝሙር አስጠኚዎች በሙሉ ምስጋናችን የላቀ ነው።"
    ], font_size=20)

    # ── SLIDE 3–4: ማውጫ ────────────────────────────────────────────
    # Split into two slides (max ~8 items each)
    add_slide("ማውጫ (Table of Contents) — 1/2", [
        "አላማ (Objectives of the Proposal)",
        "ምዕራፍ 1: መሪ / አስጠኚ",
        "  1.1. መሪነት በጥንተ ተፈጥሮ",
        "  1.2. የመሪነት ትርጉም",
        "  1.3. የመዝሙር አስጠኚነት ከመሪነት አንጻር",
        "ምዕራፍ 2: መዝሙር አስጠኚነት እስካሁን የመጣበት መንገድ",
        "  2.1. የታሪክ ጉዞና ሚና",
        "  2.2. ቀደም ሲል የነበሩ አሠራሮችና የታዩ ክፍተቶች",
    ], font_size=19)

    add_slide("ማውጫ (Table of Contents) — 2/2", [
        "ምዕራፍ 3: አዲሱ የተሻሻለ የመፍትሔ ሐሳብና የመዋቅር ፍልስፍና",
        "  3.1–3.6. Rationale | አደረጃጀት ዲያግራም | Dual Liaisons | Class Coordinator",
        "ምዕራፍ 4: የመዝሙር አስጠኚዎች አዲሱ መዋቅር",
        "  4.1. ለመዝሙር አስጠኚነት ለመመረጥ የሚያስፈልጉ መስፈርቶች",
        "  4.2–4.3. የስራ ድርሻ እና ግዴታዎች",
        "  4.4. ድርብ የመረጃ አያያዝ / Data Reconciliation",
        "ምዕራፍ 5: ተተኪ አስጠኚ | ደንቦች | ምዘና | ቅጣት | ማጠቃለያ",
    ], font_size=19)

    # ── SLIDE 5: አላማ ───────────────────────────────────────────────
    add_slide("አላማ (Objectives)", [
        "1. ወጥ የአሠራር ስርዓት መዘርጋት — መዝሙር የማስጠናትና የመማር ሂደቱን ወጥ፣ ዘመናዊና መመሪያን የተከተለ ማድረግ።",
        "2. የአስጠኚዎችን ጫና መቀነስ — ሙሉ ትኩረታቸውን በዜማ ጥራትና ይዘት ላይ ብቻ እንዲያደርጉ ማስቻል።",
        "3. ፍትሃዊነትና ቁጥጥርን ማስፈን — Checks and Balances ሥርዓት መዘርጋት።",
        "4. የመረጃ ጥራትን ማረጋገጥ — Dual Data Collection አሠራር።",
        "5. ተተኪ አስጠኚዎችን ማብቃት — በስርዓት ማፍራት።"
    ], font_size=20)

    # ── SLIDE 6: ምዕራፍ 1 — 1.1 ────────────────────────────────────
    add_slide("ምዕራፍ 1: መሪ / አስጠኚ — 1.1 መሪነት በጥንተ ተፈጥሮ", [
        "መሪነት ለሰው ልጅ ከተሰጡት አምላካዊ ጸጋዎች አንዱ ሲሆን፣ ይህም እግዚአብሔርን በመምሰል የመግዛትና የመምራት ስጦታ ነው (ዘፍ 1፥26)።",
        "እግዚአብሔር በባሕርዩ ገዢ እንደመሆኑ መጠን፣ ሰውን በጸጋውና በፍቅሩ ፍጥረታትን እንዲገዛና እንዲመራ አድርጎታል።",
        "መሪነት ለሰማያውያንም ሆነ ምድራውያን ፍጥረታት የዛሬ ሕይወታቸው መመሪያና የወደፊት ጉዟቸው አቅጣጫ እንዲሆን የተሰጠ ሥርዓት ነው።",
        "የክርስትና ሕይወት ከትውልድ ወደ ትውልድ በቀደሙት አባቶች መሪነት ሲተላለፍ የመጣ ሕይወት ነው።"
    ], font_size=19)

    # ── SLIDE 7: 1.2 ───────────────────────────────────────────────
    add_slide("1.2. የመሪነት ትርጉም", [
        "በጎ ተፅዕኖ ማሳደር፦ መሪነት ማለት በሰዎች ላይ በጎ ተፅዕኖ ማሳደር እና አርኣያ መሆን ማለት ነው።",
        "የክህሎት ዕድገት፦ ተፅዕኖ የማሳደር ክህሎት በልምድና በትምህርት ሊያድግ የሚችል ነው። አንድ መሪ ተመሪዎቹ ሀላፊነቱን በፍላጎትና በደስታ እንዲወጡ ማድረግ ይኖርበታል።",
        "መክሊትን ማብዛት፦ በክርስቶስ ክርስቲያን የሆነ ሁሉ በተሰጠው መክሊት መጠን ፍሬ ማፍራት አለበት፤ ስለሆነም እንደ ሙሴ የስጦታውን ሐሳብ ተረድቶ እንደ አሮን ድክመቱን በባልንጀራው የሚሞላበትን መንገድ መፈለግ መሪነት ነው።"
    ], font_size=20)

    # ── SLIDE 8: 1.3 (1/2) ─────────────────────────────────────────
    add_slide("1.3. የመዝሙር አስጠኚነት ከመሪነት አንጻር (1/2)", [
        "«የእውነትን ቃል በቅንነት የማናገር የማያሳፍርም ሰራተኛ ሆነህ የተፈተነውን ራስህን ለእግዚአብሔር ልታቀርብ ትጋ» (2ኛ ጢሞ 2፥15)",
        "የእውነትን ቃል የሚናገር፦ ሰማያዊ ያሬዳዊ ዜማና የእግዚአብሔርን ሕያው ቃል በቅንነት ለአባላቱ የሚያስተላልፍ አስጠኚ።",
        "የማያሳፍር፦ እውነተኛ አገልግሎት የራስን ስም ሳይሆን የክርስቶስን ክብር ብቻ የሚገልጥ ነው። አስጠኚው ለሰው ሳይሆን ለእግዚአብሔር በመታዘዝ ተልዕኮውን ይወጣል።"
    ], font_size=19)

    # ── SLIDE 9: 1.3 (2/2) ─────────────────────────────────────────
    add_slide("1.3. የመዝሙር አስጠኚነት ከመሪነት አንጻር (2/2)", [
        "ሠራተኛ፦ ለክርስቶስ አካልና ለቤተክርስቲያን ታዛዥ የሆነ፣ ለበጎ ለውጥ ራሱን የሚያዘጋጅ፣ አገልግሎቱን በትጋትና በትሕትና የሚወጣ።",
        "የተፈተነ፦ በማንኛውም ፈተና ውስጥ አልፎ ከኢየሱስ ክርስቶስ ጋር መጽናትን የተማረ፣ የሰንበት ትምህርት ቤቱን ትልቅ ሀላፊነት ለመሸከም ብቁ የሆነ አገልጋይ።"
    ], font_size=20)

    # ── SLIDE 10: ምዕራፍ 2 — 2.1 ───────────────────────────────────
    add_slide("ምዕራፍ 2: የመዝሙር አስጠኚነት እስካሁን የመጣበት መንገድ — 2.1", [
        "የታሪክ ጉዞና ሚና፦ ሰንበት ትምህርት ቤታችን ከ1965 ዓ.ም ጀምሮ እስከዛሬ ዋነኛ ሚናውን ሲወጣ የቆየው የመዝሙር ክፍል ነው።",
        "ያሬዳዊ ዜማን መጠበቅ፦ ስብሐተ እግዚአብሔርን ያሬዳዊ ዜማቸውን ሳይለቁ ከትውልድ ወደ ትውልድ በማስተላለፍ ሰንበት ትምህርት ቤቱ አሁን ላለበት ደረጃ ተደርሷል።",
        "የአስጠኚዎች ቁልፍ ድርሻ፦ በሀገረ ስብከቱ ካሉ እህት ሰንበት ትምህርት ቤቶች ሁሉ ተሻለ ውጤት ለተገኘበት ምክንያት የመዝሙር አስጠኚዎች ሚና የላቀ ነው።"
    ], font_size=19)

    # ── SLIDE 11: 2.2 ──────────────────────────────────────────────
    add_slide("2.2. ቀደም ሲል የነበሩ አሠራሮችና ክፍተቶች", [
        "የነበሩ ሥራዎች፦ ደብተር ማጻፍ፣ ዜማ ማስጠና፣ በአውደ ምሕረት ማስዘመር፣ አቴንዳንስ መያዝ፣ ዲሲፕሊን ማስከበር፣ አባላትን መምከር።",
        "ቅጣቶች (ቀደም ሲል)፦ ወጥ መመሪያ ባለመኖሩ ግርፋት፣ ወላጅ ማስመጣት፣ ልብስ ማሳጠብ፣ ከአገልግሎት ማገድ ይፈጸሙ ነበር።",
        "ቅርርብ፦ አስጠኚው አባሉን እንደ ልጁ፣ አባሉም አስጠኚውን እንደ ወላጅ ያይ ነበር፤ አባላት ሲጠፉ ቤት ድረስ ይፈለጉ ነበር።",
        "የታዩ ደካማ ጎኖች፦ ተተኪ አለማፍራት | ዝርው ማስጠናት | አስጠኚ ሲቀያየር ክፍተት | ማዳላት | ግላዊ ሕይወት መጎዳት።"
    ], font_size=19)

    # ── SLIDE 12: ምዕራፍ 3 — Rationale ────────────────────────────
    add_slide("ምዕራፍ 3: አዲሱ የተሻሻለ አቅጣጫ — 3.1 Rationale", [
        "Class Coordinator (የክፍል ተጠሪ)፦ አስጠኚዎችን ከአስተዳደራዊ ጫናዎች (አቴንዳንስ፣ ዲሲፕሊን፣ ክትትል) ነፃ ለማድረግ አዲስ ሐሳብ ሆኖ ቀርቧል።",
        "Dual Liaisons (አገናኝ ንዑሳን)፦ ሰዓትና ቦታ ሲመደብ አድልዎ እንዳይፈጠር፣ Checks and Balances ሥርዓት ለማስፈን ቀርቧል።",
        "Mezmur Instructor (አዲሱ ሚና)፦ አስጠኚው ሙሉ ትኩረቱን ያሬዳዊ ዜማን ማስተላለፍ ላይ ያደርጋል፤ ድርብ ኦዲት ዳታ የመያዝ ኃላፊነቱ ግን ይቀጥላል።"
    ], font_size=19)

    # ── SLIDE 13: 3.2 ──────────────────────────────────────────────
    add_slide("3.2. የሪፖርትና የበላይ አመራር መዋቅር (Members Affairs)", [
        "አምስቱም የክፍል ተጠሪዎች (የጎልማሳ፣ የወጣት፣ የታዳጊ፣ የሕፃናት እና የደቂቃን) ተግባራቸውን የሚያከናውኑት የአባላት ጉዳይ (Members Affairs) ክፍል ክትትል ሥር ነው።",
        "የአባላት ጉዳይ ክፍል ሚና፦ ሁሉንም ክፍላት ያስተባብራል፣ ይቆጣጠራል፣ የሥራ አፈጻጸም ይገምግማል።",
        "ውሳኔ ሰጪ አካል፦ ከክፍል ተጠሪዎች አቅም በላይ ለሆኑ ጉዳዮች (ዲሲፕሊን፣ አስተዳደር) የመጨረሻ መፍትሔ ይሰጣል።"
    ], font_size=20)

    # ── SLIDE 14: 3.3 ──────────────────────────────────────────────
    add_slide("3.3. የተሟላው የመዋቅር አደረጃጀት ዲያግራም", [
        "[ አባላት ጉዳይ — Members Affairs ]",
        "  (የአምስቱም ክፍላት ተጠሪዎች የበላይ መሪ)",
        "       |",
        "  [ጎልማሳ ተጠሪ]  [ወጣት ተጠሪ]  [ታዳጊ ተጠሪ]  [ሕፃናት ተጠሪ]  [ደቂቃን ተጠሪ]",
        "       |",
        "  ሀ. ማስተባበር (Coordination) — ለሁሉም ክፍሎች ጊዜና ቦታ ማመቻቸት",
        "  ለ. አባትነት / Leadership — ሕይወትን በቅርበት መምራት"
    ], font_size=18)

    # ── SLIDE 15: 3.4 ──────────────────────────────────────────────
    add_slide("3.4. Dual Liaisons — አገናኝ ንዑሳን", [
        "ሚናቸው፦ አገልግሎቱ ሚዛናዊ፣ ፍትሃዊ እና ጥራቱን የጠበቀ እንዲሆን ሁለት አገናኝ ንዑሳን ይመደባሉ።",
        "አመራረጥ፦",
        "  • የትምህርት ክፍል አገናኝ ንዑስ — ከትምህርት ክፍል ቀጥታ ይወከላል።",
        "  • የመዝሙር ክፍል አገናኝ ንዑስ — ከመዝሙር ክፍል ቀጥታ ይወከላል።",
        "ሁለቱም ከክፍሎቻቸው፣ ከአባላት ጉዳይ እና ከክፍል ተጠሪዎች ጋር የድልድይ ሚና ይጫወታሉ።"
    ], font_size=19)

    # ── SLIDE 16: 3.4.2 ────────────────────────────────────────────
    add_slide("3.4.2. ዋና ዋና ኃላፊነቶች — Dual Liaisons", [
        "ሀ. ፍትሃዊነት ማረጋገጥ፦ ተጠሪዎች ጊዜና ቦታ ሲመድቡ አድልዎ አለማድረጋቸውን ይቆጣጠራሉ። ካጋጠመ ወዲያውኑ ለበላይ ያሳውቃሉ።",
        "ለ. Performance Auditing፦ የወካይ ክፍሎቻቸው አገልግሎት ጥራቱን የጠበቀ፣ በዕቅድ የሚመራ መሆኑን ይሞላሉ።",
        "ሐ. ዳታ ማድረስ፦ ተጣርቶ የተዘጋጀ መረጃን ከአባላት ጉዳይ ተቀብለው ለወካዮቻቸው ያስተላልፋሉ።"
    ], font_size=19)

    # ── SLIDE 17: 3.5 ──────────────────────────────────────────────
    add_slide("3.5. የተሻሻለው አጠቃላይ የመዋቅር ዲያግራም", [
        "[ አባላት ጉዳይ — Members Affairs ]",
        "         |",
        "  [ክፍል ተጠሪዎች]  <──────>  [Dual Liaisons (ትምህርት)] & [Dual Liaisons (መዝሙር)]",
        "         |",
        "  [መዝሙር ክፍል]    [ትምህርት ክፍል]    [ኪነ-ጥበብ ክፍል]"
    ], font_size=19)

    # ── SLIDE 18: 3.6.1 — Class Coordinator Criteria ───────────────
    add_slide("3.6.1. Class Coordinator — የመመረጫ መስፈርቶች", [
        "• የአገልግሎት ቆይታ — ቢያንስ 3 ዓመት ያገለገለ",
        "• መንፈሳዊነትና ስነ-ምግባር — የተመሰከረለት፣ ለአባላቱ አርኣያ",
        "• የአመራርና ተግባቦት ክህሎት — ማስተባበር፣ ማዳመጥ፣ ቀጥታ ሠርቶ ማሳየት",
        "• አባትነት / እናትነት — ለአባላት ቅርብ የሚሆን ሥነ-ልቦናዊ ብስለት",
        "• ጊዜና ፈቃደኝነት — ኃላፊነቱን ለመወጣት ዝግጁ",
        "• የዕድሜ ገደብ — ለታዳጊ ክፍል 22+፣ ለወጣት ክፍል 25+"
    ], font_size=19)

    # ── SLIDE 19: 3.6.2 — Class Coordinator Duties ─────────────────
    add_slide("3.6.2. Class Coordinator — ዋና የሥራ ድርሻዎች", [
        "ሀ. አስተዳደርና Master Schedule፦ የክፍሉን ጊዜ ሰሌዳ ያወጣል፣ ለሁሉም ክፍሎች ሰዓትና ቦታ ያመቻቻል።",
        "ለ. የአባላት ክትትልና አባትነት፦ አጠቃላይ የአባላት መረጃ ይከታተላል፣ የጠፉትን ይፈልጋል፣ ይመክራል።",
        "ሐ. ቀጥታ ቅንጅት፦ ከመዝሙር አስጠኚው እና ከሥርዓተ ትምህርት ቁጥጥር ንዑስ ጋር ቀጥታ ይሠራል፤ ለበላይ ሪፖርት ያደርጋል።"
    ], font_size=19)

    # ── SLIDE 20: ምዕራፍ 4 — 4.1 Criteria ─────────────────────────
    add_slide("ምዕራፍ 4: የመዝሙር አስጠኚ — 4.1 የመመረጫ መስፈርቶች", [
        "• የተተኪ አስጠኚነት ልምድ ያለው/ያላት",
        "• የመዝሙር ተሰጥኦና ዕውቀት",
        "• ለማስጠናት እና ለመምራት ብቁ",
        "• ከአባላት ጋር ጥሩ ተግባቦት ያለው/ያላት",
        "• ከ2 ዓመት በላይ ያገለገለ/ች",
        "• ፈቃደኛ እና ጊዜ ያለው/ያላት",
        "• በስነ-ምግባሩ/ሯ የተመሰከረለት/ላት",
        "• የታዳጊ፣ ወጣት ወይም ጎልማሳ ክፍል አባል"
    ], font_size=19)

    # ── SLIDE 21: 4.2 — Duties ─────────────────────────────────────
    add_slide("4.2. የስራ ድርሻ — Mezmur Instructor", [
        "• መዝሙሮች በክፍሉ እንዲፃፉ ያደርጋል፣ ዜማቸውን ያስጠናል",
        "• ሳምንታዊ የአስጠኚዎች ጥናት ያካፍላል",
        "• አዳዲስ መዝሙሮችን አፅፎ ዜማቸውን ያስጠናል",
        "• የበዓላት ዝማሬ ያስጠናል",
        "• ሽብሸባ፣ እንቅስቃሴ፣ ከበሮና ወረብ ያስጠናል",
        "• አውደ ምሕረት ላይ ዜማ እንዲዘመር ያደርጋል"
    ], font_size=19)

    # ── SLIDE 22: 4.3 — Obligations ────────────────────────────────
    add_slide("4.3. ግዴታዎች — Mezmur Instructor", [
        "• መመሪያ ተግብሮ ማስተግበር",
        "• ሰዓቱን አክብሮ ቀጥሎ ሃላፊነቱን መወጣት",
        "• የአስጠኚዎች ጥናት ላይ ሁልጊዜ መገኘት",
        "• ተተኪ አስጠኚዎችን ማብቃት",
        "• ሳይያስፈቅድ ሳይተካ አለመቅረት",
        "• ከመአንክ እና ከመክተ ያልሆነ መመሪያ አለማስፈጸም"
    ], font_size=19)

    # ── SLIDE 23: 4.4 ──────────────────────────────────────────────
    add_slide("4.4. ከClass Coordinator ጋር ያለው ቅንጅት", [
        "ጊዜ ሰሌዳ፦ የሰዓትና ቦታ ጥያቄ ለክፍል ተጠሪው ቀጥታ ይቀርባል።",
        "ዲሲፕሊን፦ በጥናት ሰዓት ጸጥታ የማስከበር ኃላፊነት የአስጠኚው ነው፤ ተጠሪው ከርቀት ያዛል።",
        "Dual Data Collection፦ ተጠሪው አጠቃላይ የአባላት መረጃ ይይዛል፤ አስጠኚው የጥናት ሰዓት ተሳትፎ ይከታተላል።",
        "Data Reconciliation፦ ሁለቱ ሲሰሜሙ ትክክለኛ ምስል ይቀርጻሉ። ክፍተት ካለ ወዲያው ይታወቃል።"
    ], font_size=19)

    # ── SLIDE 24: ምዕራፍ 5 — 5.1.1 ────────────────────────────────
    add_slide("ምዕራፍ 5: የተተኪ አስጠኚ — 5.1.1 የመመረጫ መስፈርቶች", [
        "• ከ2 ዓመት ያላነሰ የአገልግሎት ቆይታ",
        "• የመዝሙር ተሰጥኦ (ዜማ፣ ዝማሬ) እና ዕውቀት",
        "• ለማስጠናት ፍላጎትና አቅም ያለው/ያላት",
        "• ከአባላት ጋር መልካም ተግባቦት",
        "• በስነ-ምግባሩ/ሯ የተመሰከረለት/ላት",
        "• የታዳጊ፣ ወጣት ወይም ጎልማሳ ክፍል አባል"
    ], font_size=19)

    # ── SLIDE 25: 5.1.2 — Sub Duties ──────────────────────────────
    add_slide("5.1.2. ተተኪ አስጠኚ — ዋና ዋና ሥራዎች", [
        "• አስጠኚ ሲቀር ወዲያው ወስዶ ያስጠናቸዋል",
        "• መዝሙር እንዲፅፉ ያደርጋል",
        "• የአስጠኚዎችን ጥናት ያካፍላል",
        "• ከአስጠኚው ጋር ጥናት ያስጠናቸዋል፣ ክህሎታቸውን ያዳምጣል"
    ], font_size=20)

    # ── SLIDE 26: 5.2 — Ethics & Prohibitions ──────────────────────
    add_slide("5.2. ስነ-ምግባር ደንቦች እና ክልከላዎች", [
        "5.2.1. የስነ-ምግባር ደንቦች፦ መንፈሳዊ፣ ትሁት፣ ታዛዥ፣ አርኣያ የሚሆን ሕይወት። ምስጢር ጠባቂ።",
        "5.2.2. ተከልክሎ ያሉ ነገሮች፦",
        "  ✗ ከስልጣን ወሰን ውጭ — አባላትን መቅጣት፣ ማገድ ወይም ማባረር",
        "  ✗ አላስፈላጊ ቅርርብ — ማዳላት ወይም ከስምሪቱ ውጭ ግንኙነት",
        "  ✗ ያልተፈቀዱ ዝማሬዎች — ያልተገምገሙ መዝሙሮችን ማስጠና",
        "  ✗ ያለፈቃድ መቅረት — ሳያስፈቅዱ ሳይተኩ አለመቅረብ"
    ], font_size=19)

    # ── SLIDE 27: 5.3 — Assessment ─────────────────────────────────
    add_slide("5.3. የምዘና ስልቶች", [
        "• የቃል ምዘና",
        "• የጽሑፍ ምዘና",
        "• የተግባር / የአውደ ምሕረት ምዘና",
        "• የክትትል ምዘና (Monitoring)"
    ], font_size=22)

    # ── SLIDE 28: 5.4 — Penalties ──────────────────────────────────
    add_slide("5.4. የቅጣትና የእርምት እርምጃዎች", [
        "1ኛ ደረጃ — የቃል ማስጠንቀቂያ፦ ያለምክንያት 1 ጊዜ ሲቀር ወይም አነስተኛ ግድፈት ሲፈጽም",
        "2ኛ ደረጃ — የጽሑፍ ማስጠንቀቂያ፦ ጥፋቱን ሲደግም ወይም ሳያስፈቅድ ከ2 ጊዜ በላይ ሲቀር",
        "3ኛ ደረጃ — ጊዜያዊ ማገድ፦ መመሪያ ለማስፈጸም ፈቃደኛ ሳይሆን ሲቀር ወይም ከባድ ግድፈት ሲፈጽም",
        "4ኛ ደረጃ — ሙሉ ማሰናበት፦ ጥፋቱ ከባድ ሆኖ ወይም የሰንበት ትምህርት ቤቱን ክብር የሚጎዳ ሆኖ ሲገኝ"
    ], font_size=19)

    # ── SLIDE 29: 5.5 — Conclusion ─────────────────────────────────
    add_slide("5.5. ማጠቃለያ (Conclusion)", [
        "«ሰውነታችሁን እግዚአብሔርን ደስ የሚያሰኝና ሕያው ቅዱስም መሥዋዕት አድርጋችሁ ታቀርቡ ዘንድ በእግዚአብሔር ርኅራኄ እለምናችኋለሁ፤ እርሱም ለአእምሮ የሚመች አገልግሎታችሁ ነው» (ሮሜ 12፥1)",
        "ይህ አዲስ መዋቅር ከተተገበረ፦",
        "  ✓ የአስጠኚዎች ጫና ይቀንሳል",
        "  ✓ የአባላት ክትትልና አባትነት በጥራት ይከናወናል",
        "  ✓ የመረጃ ጥራትና ፍትሃዊነት ይረጋገጣል",
        "  ✓ ብቁ ተተኪ አስጠኚዎች ይፈራሉ",
        "",
        "ወስበሐት ለእግዚአብሔር!"
    ], font_size=20)

    prs.save('Mezmur_Structure_Proposal.pptx')
    print(f"Successfully created Mezmur_Structure_Proposal.pptx with {len(prs.slides)} slides")

if __name__ == '__main__':
    main()
