import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# New Brand Colors (Unique and Professional)
DEEP_TEAL = RGBColor(0x00, 0x33, 0x4D)
GOLDEN_ACCENT = RGBColor(0xD4, 0xAF, 0x37)
SOFT_WHITE = RGBColor(0xF5, 0xF5, 0xF5)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)

def apply_background(slide, color=DEEP_TEAL):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(prs, title_text, subtitle_text=""):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    apply_background(slide, DEEP_TEAL)
    
    # Decorator line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.5), Inches(8), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = GOLDEN_ACCENT
    line.line.color.rgb = GOLDEN_ACCENT
    
    title = slide.shapes.title
    title.text = title_text
    title.top = Inches(2.5)
    title.left = Inches(0.5)
    title.width = Inches(9.0)
    for paragraph in title.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.color.rgb = GOLDEN_ACCENT
            run.font.bold = True
            run.font.name = "Nyala" # Good Amharic font if available, fallback to Arial
            run.font.size = Pt(44)
            
    if subtitle_text:
        subtitle = slide.shapes.placeholders[1]
        subtitle.text = subtitle_text
        subtitle.top = Inches(4.0)
        subtitle.left = Inches(0.5)
        subtitle.width = Inches(9.0)
        for paragraph in subtitle.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.color.rgb = SOFT_WHITE
                run.font.bold = False
                run.font.size = Pt(28)
    else:
        if len(slide.shapes.placeholders) > 1:
            sp = slide.shapes.placeholders[1]
            sp.element.getparent().remove(sp.element)
                
    return slide

def add_divider_slide(prs, title_text):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    apply_background(slide, GOLDEN_ACCENT)
    title = slide.shapes.title
    title.text = title_text
    title.top = Inches(3.0)
    title.left = Inches(0.5)
    title.width = Inches(9.0)
    for paragraph in title.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.color.rgb = DEEP_TEAL
            run.font.bold = True
            run.font.size = Pt(50)
            
    # Remove unused subtitle placeholder
    if len(slide.shapes.placeholders) > 1:
        sp = slide.shapes.placeholders[1]
        sp.element.getparent().remove(sp.element)
        
    return slide

def add_content_slide(prs, title_text, body_text):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_background(slide, SOFT_WHITE)
    
    # Top Banner
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.2))
    banner.fill.solid()
    banner.fill.fore_color.rgb = DEEP_TEAL
    banner.line.color.rgb = DEEP_TEAL
    
    title_shape = slide.shapes.title
    title_shape.text = title_text
    title_shape.left = Inches(0.5)
    title_shape.top = Inches(0.2)
    title_shape.width = Inches(9.0)
    title_shape.height = Inches(0.8)
    for paragraph in title_shape.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT
        for run in paragraph.runs:
            run.font.color.rgb = GOLDEN_ACCENT
            run.font.bold = True
            run.font.size = Pt(32)
            
    tf = slide.shapes.placeholders[1].text_frame
    slide.shapes.placeholders[1].top = Inches(1.5)
    slide.shapes.placeholders[1].left = Inches(0.5)
    slide.shapes.placeholders[1].width = Inches(9.0)
    slide.shapes.placeholders[1].height = Inches(5.5)
    
    lines = body_text.strip().split('\n')
    for i, line in enumerate(lines):
        line = line.strip()
        if not line:
            continue
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.space_after = Pt(10)
        for run in p.runs:
            run.font.color.rgb = DARK_GRAY
            run.font.size = Pt(22)
            
    return slide

def parse_extracted_text(filepath):
    slides = []
    current_slide = []
    
    with open(filepath, 'r', encoding='utf-8') as f:
        for line in f:
            if line.startswith('--- Slide'):
                if current_slide:
                    slides.append('\n'.join(current_slide).strip())
                current_slide = []
            else:
                current_slide.append(line.strip())
        if current_slide:
            slides.append('\n'.join(current_slide).strip())
            
    return slides

def process_slides(prs, raw_slides):
    for raw_slide in raw_slides:
        if not raw_slide.strip():
            continue
        lines = raw_slide.split('\n')
        title = lines[0]
        body = '\n'.join(lines[1:]) if len(lines) > 1 else ""
        if len(title) > 60:
            body = title + "\n" + body
            title = "ማስታወሻ" # Note
        add_content_slide(prs, title, body)

def main():
    base_dir = "/home/dawit/.gemini/antigravity-ide/brain/996f3f23-b2a3-4157-95ad-f7bb11cf9ec5/scratch"
    combo_file = os.path.join(base_dir, "extracted_text.txt")
    hawariyaw_file = os.path.join(base_dir, "hawariyaw_extracted.txt")
    homiletics_file = os.path.join(base_dir, "homiletics_extracted.txt")
    
    combo_slides = parse_extracted_text(combo_file)
    hawariyaw_slides = parse_extracted_text(hawariyaw_file)
    homiletics_slides = parse_extracted_text(homiletics_file)
    
    # --- Presentation 1: Homiletics ---
    prs1 = Presentation()
    add_title_slide(prs1, "የስብከት ዘዴ", "Homiletics (Preaching Methods)")
    add_divider_slide(prs1, "ክፍል 1: ዋና ትምህርት")
    # Combo slides 3 to 43 (0-indexed 2 to 42)
    process_slides(prs1, combo_slides[2:43])
    add_divider_slide(prs1, "ክፍል 2: ተጨማሪ ጥቅሶችና መረጃዎች")
    process_slides(prs1, homiletics_slides)
    
    output1 = "/home/dawit/projects/amdehaymanot official website/hawariyaw teliko/የስብከት_ዘዴ_New.pptx"
    prs1.save(output1)
    print(f"Created {output1} with {len(prs1.slides)} slides.")
    
    # --- Presentation 2: Hawariyaw Teliko ---
    prs2 = Presentation()
    add_title_slide(prs2, "ሐዋርያዊ ተልዕኮ", "Apostolic Mission")
    add_divider_slide(prs2, "ክፍል 1: ዋና ትምህርት")
    # Combo slides 44 to 51 (0-indexed 43 to 50)
    process_slides(prs2, combo_slides[43:51])
    add_divider_slide(prs2, "ክፍል 2: ተጨማሪ ጥቅሶችና መረጃዎች")
    process_slides(prs2, hawariyaw_slides)
    
    output2 = "/home/dawit/projects/amdehaymanot official website/hawariyaw teliko/ሐዋርያዊ_ተልዕኮ_New.pptx"
    prs2.save(output2)
    print(f"Created {output2} with {len(prs2.slides)} slides.")

if __name__ == "__main__":
    main()
