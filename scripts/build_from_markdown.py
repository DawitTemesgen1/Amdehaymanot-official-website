import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# Premium Orthodox Palette
MIDNIGHT_TEAL = RGBColor(0x00, 0x2B, 0x40)  # #002B40
DEEP_TEAL     = RGBColor(0x00, 0x33, 0x4D)  # #00334D
GOLDEN_ACCENT = RGBColor(0xD4, 0xAF, 0x37)  # #D4AF37
GOLD_LIGHT    = RGBColor(0xF0, 0xD3, 0x76)  # #F0D376
CANVAS_BG     = RGBColor(0xF4, 0xF6, 0xF9)  # #F4F6F9
CARD_BORDER   = RGBColor(0xD2, 0xDC, 0xE5)  # #D2DCE5
TEXT_DARK     = RGBColor(0x1A, 0x20, 0x2C)  # #1A202C
TEXT_BODY     = RGBColor(0x2D, 0x37, 0x48)  # #2D3748
SCRIPTURE_CLR = RGBColor(0x00, 0x4D, 0x66)  # #004D66
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)

def apply_solid_fill(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def create_title_slide(prs, title, subtitle, scripture="", course_name="የስብከት ዘዴ"):
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # blank layout
    apply_solid_fill(slide, MIDNIGHT_TEAL)
    
    # Top Gold Border Accent
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.14))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = GOLDEN_ACCENT
    top_bar.line.fill.background()
    
    # Bottom Gold Border Accent
    bot_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.36), Inches(10), Inches(0.14))
    bot_bar.fill.solid()
    bot_bar.fill.fore_color.rgb = GOLDEN_ACCENT
    bot_bar.line.fill.background()
    
    # Decorative Inner Frame Line
    frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0.4), Inches(9.2), Inches(6.7))
    frame.fill.background()
    frame.line.color.rgb = RGBColor(0x78, 0x62, 0x23)
    frame.line.width = Pt(1.0)
    
    clean_title = re.sub(r'[\*\#]', '', title).strip()
    clean_subtitle = re.sub(r'[\*\#]', '', subtitle).strip() if subtitle else ""
    clean_scripture = re.sub(r'[\*\#]', '', scripture).strip() if scripture else ""
    
    # Top Ornamental Pill Badge
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.8), Inches(0.85), Inches(4.4), Inches(0.55))
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor(0x06, 0x38, 0x52)
    badge.line.color.rgb = GOLDEN_ACCENT
    badge.line.width = Pt(1.2)
    btf = badge.text_frame
    btf.word_wrap = False
    btf.margin_top = Inches(0.08)
    bp = btf.paragraphs[0]
    bp.text = f"❖  {course_name}  ❖"
    bp.alignment = PP_ALIGN.CENTER
    for run in bp.runs:
        run.font.color.rgb = GOLD_LIGHT
        run.font.size = Pt(17)
        run.font.bold = True
        run.font.name = "Segoe UI"
    
    # Main Title Box
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(2.2), Inches(8.8), Inches(1.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = clean_title
    p.alignment = PP_ALIGN.CENTER
    for run in p.runs:
        run.font.color.rgb = GOLDEN_ACCENT
        run.font.bold = True
        run.font.size = Pt(46)
        run.font.name = "Segoe UI"
        
    # Subtitle / Attribution Box
    if clean_subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.7), Inches(8.4), Inches(2.0))
        stf = sub_box.text_frame
        stf.word_wrap = True
        sp = stf.paragraphs[0]
        sp.text = clean_subtitle
        sp.alignment = PP_ALIGN.CENTER
        for run in sp.runs:
            run.font.color.rgb = WHITE
            run.font.size = Pt(21)
            run.font.name = "Segoe UI"
            
    # Scripture Callout Card at Bottom (if present)
    if clean_scripture:
        quote_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.4), Inches(8.4), Inches(2.3))
        quote_box.fill.solid()
        quote_box.fill.fore_color.rgb = RGBColor(0x04, 0x30, 0x47)
        quote_box.line.color.rgb = GOLDEN_ACCENT
        quote_box.line.width = Pt(1.5)
        qtf = quote_box.text_frame
        qtf.word_wrap = True
        qtf.margin_left = Inches(0.5)
        qtf.margin_right = Inches(0.5)
        qtf.margin_top = Inches(0.28)
        qtf.margin_bottom = Inches(0.2)
        qp = qtf.paragraphs[0]
        qp.text = clean_scripture
        qp.alignment = PP_ALIGN.CENTER
        for run in qp.runs:
            run.font.color.rgb = GOLD_LIGHT
            run.font.italic = True
            run.font.size = Pt(18.5)
            run.font.name = "Segoe UI"

def create_divider_slide(prs, section_title, course_name="የስብከት ዘዴ"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_solid_fill(slide, MIDNIGHT_TEAL)
    
    # Top & Bottom Gold Trims
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.12))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = GOLDEN_ACCENT
    top_bar.line.fill.background()
    
    bot_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.38), Inches(10), Inches(0.12))
    bot_bar.fill.solid()
    bot_bar.fill.fore_color.rgb = GOLDEN_ACCENT
    bot_bar.line.fill.background()
    
    # Center Floating Regal Card
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.8), Inches(8.0), Inches(3.8))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x04, 0x34, 0x4D)
    card.line.color.rgb = GOLDEN_ACCENT
    card.line.width = Pt(2.0)
    
    ctf = card.text_frame
    ctf.word_wrap = True
    ctf.margin_left = Inches(0.5)
    ctf.margin_right = Inches(0.5)
    ctf.margin_top = Inches(0.65)
    
    # Course Name Top Header
    p_dec = ctf.paragraphs[0]
    p_dec.text = f"❖  {course_name}  ❖"
    p_dec.alignment = PP_ALIGN.CENTER
    for run in p_dec.runs:
        run.font.color.rgb = GOLD_LIGHT
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.name = "Segoe UI"
        
    # Main Section Header (Single clean entry)
    clean_section = re.sub(r'[\*\#]', '', section_title).strip()
    p_main = ctf.add_paragraph()
    p_main.text = clean_section
    p_main.space_before = Pt(24)
    p_main.alignment = PP_ALIGN.CENTER
    for run in p_main.runs:
        run.font.color.rgb = GOLDEN_ACCENT
        run.font.bold = True
        run.font.size = Pt(38)
        run.font.name = "Segoe UI"

def create_toc_slide(prs, slide_title, items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_solid_fill(slide, CANVAS_BG)
    
    # Top Header Banner
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.18))
    header.fill.solid()
    header.fill.fore_color.rgb = MIDNIGHT_TEAL
    header.line.fill.background()
    
    # Gold accent line under header
    gold_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.18), Inches(10), Inches(0.04))
    gold_line.fill.solid()
    gold_line.fill.fore_color.rgb = GOLDEN_ACCENT
    gold_line.line.fill.background()
    
    # Header Left Accent Pillar
    pillar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.42), Inches(0.24), Inches(0.08), Inches(0.68))
    pillar.fill.solid()
    pillar.fill.fore_color.rgb = GOLDEN_ACCENT
    pillar.line.fill.background()
    
    # Title Text in Header
    title_box = slide.shapes.add_textbox(Inches(0.62), Inches(0.12), Inches(8.9), Inches(0.92))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "ማውጫ / ይዘት"
    p.alignment = PP_ALIGN.LEFT
    for run in p.runs:
        run.font.color.rgb = GOLDEN_ACCENT
        run.font.bold = True
        run.font.size = Pt(28)
        run.font.name = "Segoe UI"
        
    valid_items = [it.strip() for it in items if it.strip()]
    cleaned_items = []
    for raw in valid_items:
        t = raw
        if t.startswith("- "):
            t = t[2:].strip()
        t = t.replace("**", "").strip()
        if t:
            cleaned_items.append(t)
            
    n = len(cleaned_items)
    if n == 0:
        return
        
    geez_nums = ["፩", "፪", "፫", "፬", "፭", "፮", "፯", "፰", "፱", "፲", "፲፩", "፲፪"]
    
    start_top = 1.45
    avail_height = 5.5
    row_height = min(0.68, (avail_height - (n - 1) * 0.12) / n)
    gap = 0.12 if n <= 7 else 0.08
    
    for i, item_text in enumerate(cleaned_items):
        top_pos = start_top + i * (row_height + gap)
        
        badge_label = geez_nums[i] if i < len(geez_nums) else f"{i+1}"
        display_text = item_text
        m = re.match(r'^([፩፪፫፬፭፮፯፰፱፲\d]+)[\.\s:፡\-]*(.*)', item_text)
        if m:
            badge_label = m.group(1)
            display_text = m.group(2).strip()
            
        # Outer Card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(top_pos), Inches(8.9), Inches(row_height))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = CARD_BORDER
        card.line.width = Pt(1.2)
        
        # Badge Pill Shape
        badge_w = 0.54
        badge_h = row_height - 0.16
        badge_top = top_pos + 0.08
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.70), Inches(badge_top), Inches(badge_w), Inches(badge_h))
        badge.fill.solid()
        badge.fill.fore_color.rgb = MIDNIGHT_TEAL
        badge.line.color.rgb = GOLDEN_ACCENT
        badge.line.width = Pt(1.2)
        
        btf = badge.text_frame
        btf.word_wrap = False
        btf.margin_left = Inches(0)
        btf.margin_right = Inches(0)
        btf.margin_top = Inches(0.04)
        btf.margin_bottom = Inches(0)
        bp = btf.paragraphs[0]
        bp.alignment = PP_ALIGN.CENTER
        brun = bp.add_run()
        brun.text = badge_label
        brun.font.name = "Segoe UI"
        brun.font.size = Pt(17)
        brun.font.bold = True
        brun.font.color.rgb = GOLDEN_ACCENT
        
        # Text Label Box
        txt_left = 1.38
        txt_w = 7.9
        txt_box = slide.shapes.add_textbox(Inches(txt_left), Inches(top_pos + 0.04), Inches(txt_w), Inches(row_height - 0.08))
        ttf = txt_box.text_frame
        ttf.word_wrap = True
        ttf.margin_left = Inches(0.05)
        ttf.margin_right = Inches(0.1)
        ttf.margin_top = Inches(0.06)
        ttf.margin_bottom = Inches(0)
        tp = ttf.paragraphs[0]
        tp.alignment = PP_ALIGN.LEFT
        trun = tp.add_run()
        trun.text = display_text
        trun.font.name = "Segoe UI"
        trun.font.size = Pt(18.5)
        trun.font.bold = True
        trun.font.color.rgb = DEEP_TEAL

def create_content_slide(prs, slide_title, items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_solid_fill(slide, CANVAS_BG)
    
    # Top Header Banner
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.18))
    header.fill.solid()
    header.fill.fore_color.rgb = MIDNIGHT_TEAL
    header.line.fill.background()
    
    # Gold accent line under header
    gold_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.18), Inches(10), Inches(0.04))
    gold_line.fill.solid()
    gold_line.fill.fore_color.rgb = GOLDEN_ACCENT
    gold_line.line.fill.background()
    
    # Header Left Accent Pillar (Visual Anchor)
    pillar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.42), Inches(0.24), Inches(0.08), Inches(0.68))
    pillar.fill.solid()
    pillar.fill.fore_color.rgb = GOLDEN_ACCENT
    pillar.line.fill.background()
    
    # Clean Title Text
    clean_title = re.sub(r'[\*\#]', '', slide_title).strip()
    clean_title = re.sub(r'\s*[\—\-]?\s*ክፍል\s*\d+\s*', ' ', clean_title)
    clean_title = re.sub(r'\s*\(\s*Table of Contents.*?\)', '', clean_title, flags=re.IGNORECASE).strip()
    
    # Title Text in Header
    title_box = slide.shapes.add_textbox(Inches(0.62), Inches(0.12), Inches(8.9), Inches(0.92))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = clean_title
    p.alignment = PP_ALIGN.LEFT
    for run in p.runs:
        run.font.color.rgb = GOLDEN_ACCENT
        run.font.bold = True
        run.font.size = Pt(27)
        run.font.name = "Segoe UI"
        
    # Main Content Card (Clean white canvas card)
    content_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.40), Inches(1.36), Inches(9.2), Inches(5.82))
    content_card.fill.solid()
    content_card.fill.fore_color.rgb = WHITE
    content_card.line.color.rgb = CARD_BORDER
    content_card.line.width = Pt(1.2)
    
    ctf = content_card.text_frame
    ctf.word_wrap = True
    ctf.margin_left = Inches(0.45)
    ctf.margin_right = Inches(0.45)
    ctf.margin_top = Inches(0.38)
    ctf.margin_bottom = Inches(0.35)
    
    # Filter valid items
    valid_items = [it for it in items if it.strip()]
    num_items = len(valid_items)
    
    # Dynamic typography scaling based on item count (Extra large, zoomed presentation view)
    if num_items >= 9:
        sz_l0, sz_l1 = 20.0, 18.0
        sp_before = Pt(10)
        sp_after = Pt(5)
    elif num_items >= 6:
        sz_l0, sz_l1 = 22.0, 20.0
        sp_before = Pt(13)
        sp_after = Pt(7)
    elif num_items >= 4:
        sz_l0, sz_l1 = 24.0, 22.0
        sp_before = Pt(18)
        sp_after = Pt(10)
    else:
        sz_l0, sz_l1 = 26.5, 24.0
        sp_before = Pt(24)
        sp_after = Pt(14)
        
    is_first = True
    for raw_item in valid_items:
        indent_level = 1 if raw_item.startswith("  ") or raw_item.startswith("\t") else 0
        item_text = raw_item.strip()
        if item_text.startswith("- "):
            item_text = item_text[2:].strip()
            
        if is_first:
            p = ctf.paragraphs[0]
            is_first = False
        else:
            p = ctf.add_paragraph()
            
        p.alignment = PP_ALIGN.LEFT
        p.level = 0 # Keep all at level 0 for crisp PowerPoint rendering
        
        # Check if this is a quote line
        is_quote = "«" in item_text or "”" in item_text or item_text.startswith('"')
        
        if indent_level == 0:
            p.space_before = sp_before
            p.space_after = Pt(3)
        else:
            p.space_before = Pt(2)
            p.space_after = sp_after
            
        # Format text runs
        if indent_level == 1:
            # Add elegant golden bullet symbol
            bullet_run = p.add_run()
            bullet_run.text = "   •  " if not (item_text[0].isdigit() and item_text[1:3] in ['. ', '፣ ']) else "   "
            bullet_run.font.name = "Segoe UI"
            bullet_run.font.bold = True
            bullet_run.font.color.rgb = GOLDEN_ACCENT
            bullet_run.font.size = Pt(sz_l1)
            
        parts = re.split(r'(\*\*.*?\*\*)', item_text)
        for part in parts:
            if not part:
                continue
            run = p.add_run()
            run.font.name = "Segoe UI"
            if part.startswith('**') and part.endswith('**'):
                run.text = part[2:-2]
                run.font.bold = True
                run.font.color.rgb = DEEP_TEAL if indent_level == 0 else TEXT_DARK
                run.font.size = Pt(sz_l0) if indent_level == 0 else Pt(sz_l1)
            else:
                clean_part = part.replace('*', '')
                run.text = clean_part
                run.font.bold = False
                if is_quote:
                    run.font.italic = True
                    run.font.color.rgb = SCRIPTURE_CLR
                else:
                    run.font.color.rgb = TEXT_DARK if indent_level == 0 else TEXT_BODY
                run.font.size = Pt(sz_l0) if indent_level == 0 else Pt(sz_l1)

def parse_markdown_to_ppt(md_path, output_pptx):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    course_name = "የስብከት ዘዴ" if "YESIBIKET" in md_path.upper() else "ሐዋርያዊ ተልዕኮ"
    
    with open(md_path, 'r', encoding='utf-8') as f:
        content = f.read()
        
    sections = content.split("### ስላይድ")
    
    for sec in sections[1:]:
        lines = [l for l in sec.strip().split('\n') if l.strip()]
        if not lines:
            continue
            
        title_line = lines[0]
        match = re.search(r'^\s*\d+[\s:፡\-]*(.*)', title_line)
        slide_title = match.group(1).strip() if match else title_line
        
        filtered_body = []
        for bl in lines[1:]:
            s_bl = bl.strip()
            if not s_bl:
                continue
            if s_bl.startswith("## ") or s_bl.startswith("# ") or s_bl.startswith("---") or s_bl == "---":
                continue
            if "ሞጁል" in s_bl and ("##" in bl or "📌" in bl or "Module" in bl):
                continue
            filtered_body.append(bl)
            
        body_lines = filtered_body
        
        if "የርዕስ ስላይድ" in title_line or "Title Slide" in title_line:
            main_t = ""
            sub_t = ""
            scripture = ""
            for bl in body_lines:
                if "ዋና ርዕስ:" in bl or "ዋና ርዕስ፦" in bl:
                    main_t = bl.split(":", 1)[-1].split("፦", 1)[-1].strip()
                elif "ንዑስ ርዕስ:" in bl or "ንዑስ ርዕስ፦" in bl:
                    sub_t = bl.split(":", 1)[-1].split("፦", 1)[-1].strip()
                elif "መሪ ጥቅስ:" in bl or "መሪ ጥቅስ፦" in bl:
                    scripture = bl.split(":", 1)[-1].split("፦", 1)[-1].strip()
            create_title_slide(prs, main_t or slide_title, sub_t, scripture, course_name)
            
        elif "ማውጫ" in title_line or "Table of Contents" in title_line or "ይዘት" in title_line:
            create_toc_slide(prs, slide_title, body_lines)
            
        elif "መለያ ስላይድ" in title_line or "Section Divider" in title_line:
            divider_text = slide_title
            for bl in body_lines:
                if "ክፍል" in bl:
                    divider_text = bl.replace("-", "").replace("**", "").strip()
            create_divider_slide(prs, divider_text, course_name)
            
        else:
            create_content_slide(prs, slide_title, body_lines)
            
    prs.save(output_pptx)
    print(f"Saved {output_pptx} with {len(prs.slides)} slides.")

def main():
    base_dir = "/home/dawit/projects/amdehaymanot official website/hawariyaw teliko"
    
    f1 = os.path.join(base_dir, "YESIBIKET_ZEDE_FULL_CONTENT.md")
    out1 = os.path.join(base_dir, "የስብከት_ዘዴ_Master.pptx")
    if os.path.exists(f1):
        parse_markdown_to_ppt(f1, out1)
        
    f2 = os.path.join(base_dir, "HAWARIYAW_TELIKO_FULL_CONTENT.md")
    out2 = os.path.join(base_dir, "ሐዋርያዊ_ተልዕኮ_Master.pptx")
    if os.path.exists(f2):
        parse_markdown_to_ppt(f2, out2)

if __name__ == "__main__":
    main()
